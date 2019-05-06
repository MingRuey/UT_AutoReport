import sys
import os
from datetime import datetime as dt
import pathlib
import yaml
import pptx
from pptx.util import Inches as Inch
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR as Line

_path = os.path.dirname(os.path.dirname(__file__))
if _path not in sys.path:
    sys.path.append(_path)

from units.cover import ReportCover  # noqa E402
from units.figures import Figure  # noqa E402
from units.subjects import SubjectTitle, Text  # noqa E402

A4 = (Inch(7.5), Inch(10.83))
SLDBLANK = 6
PROJECT_DIR = pathlib.Path(os.path.realpath(__file__)).parents[1]


class Subject:

    def __init__(self, title, info, sections):
        self._title = str(title)
        self._info = str(info)

        self._sections = []
        for section in sections:
            name, items = next(iter(section.items()))
            items = [next(iter(item.keys())) for item in items]
            if "text" not in items and "picture" not in items:
                msg = "Section {} must have text and/or picture"
                raise ValueError(msg.format(name))

            self.sections.append(section)

    def __str__(self):
        display = "title: {}, info: {}, sections: {}"
        return display.format(self.title, self.info, self.sections)

    @property
    def title(self):
        return self._title

    @property
    def info(self):
        return self._info

    @property
    def sections(self):
        return self._sections

    @classmethod
    def to_yaml(cls, yml):
        raise NotImplementedError()

    @classmethod
    def from_yaml(cls, yml, path=None):
        if path is not None:
            yml = pathlib.Path(path).joinpath(yml)

        with open(str(yml), "r") as f:
            subject = yaml.safe_load(f)

        return cls(
            subject['title'],
            subject['info'],
            subject['sections']
        )


class ReportSettings():

    formats = ["WeeklyReport"]

    def __init__(
            self,
            form,
            subjects,
            author=None,
            date=None
            ):

        if form not in self.formats:
            msg = "Not supported form: {}; Can only be one of {}"
            raise NotImplementedError(msg.format(form, self.formats))

        self._formats = form
        self._author = str(author) if author else "UT-AUTO-REPORT"
        self._date = dt.strptime(date, "%Y-%m-%d").date() if date else dt.now()

        self._subject_files = []
        self._subject_objs = []
        for subject in subjects:
            self._subject_objs.append(Subject.from_yaml(subject))
            self._subject_files.append(subject)

    @property
    def title(self):
        return self._formats

    @property
    def subjects(self):
        return self._subject_objs

    @property
    def author(self):
        return self._author

    @property
    def date(self):
        return self._date

    @classmethod
    def to_yaml(cls, yml):
        raise NotImplementedError()

    @classmethod
    def from_yaml(cls, yml, path=None):
        if path is not None:
            yml = pathlib.Path(path).joinpath(yml)
            path = yml.parents[0]

        with open(str(yml), "r") as f:
            settings = yaml.safe_load(f)

        subjects = settings['subjects']
        if path is not None:
            subjects = [str(path.joinpath(sub)) for sub in subjects]

        return cls(
            form=settings['format'],
            subjects=subjects,
            date=settings['date'],
            author=settings['author']
        )


class UTSimple:

    version = ["0.0beta"]

    def __init__(self, settings: ReportSettings):

        if not isinstance(settings, ReportSettings):
            msg = "setting must be instance of ReportSettings"
            raise TypeError(msg)

        self.setting = settings

    def to_pptx(self, file):
        file = str(file)
        if not file.endswith('.pptx'):
            raise ValueError("Invalid save out file name")

        # presentation wise settings
        prs = pptx.Presentation()
        prs.slide_height = A4[0]
        prs.slide_width = A4[1]

        core = prs.core_properties
        core.author = self.setting.author
        core.created = dt.now()
        core.last_modified_by = self.setting.author
        core.last_printed = dt.now()
        core.modified = dt.now()
        core.title = self.setting.formats
        core.version = UTSimple.version

        self.prs = prs
        self._add_cover_slide()
        for subject in self.setting.subjects:
            self._add_subject_slides(subject)
        self.prs.save(file)

    def _add_cover_slide(self):

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[SLDBLANK])
        shapes = slide.shapes

        shapes.add_picture(
            str(PROJECT_DIR.joinpath('./data/banner_utechzone_blue.png')),
            left=0, top=0,
            width=self.prs.slide_width
            )

        cover = ReportCover(
            header=self.setting.title,
            titles=[sub.title for sub in self.setting.subjects],
            author=self.setting.author,
            date=self.setting.date
            )

        cover.add_to_shapes(
            shapes,
            left=1.85,
            top=1.72
            )

        shapes.add_picture(
            str(PROJECT_DIR.joinpath('./data/banner_utechzone_blue.png')),
            left=0, top=self.prs.slide_height - Inch(1.07),
            width=self.prs.slide_width
            )

    def _add_subject_slides(self, subject):

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[SLDBLANK])
        shapes = slide.shapes

        title = SubjectTitle(
            title=subject.title,
            description=subject.info
        )

        title.add_to_shapes(
            shapes,
            left=0.25, top=0.18
        )

        _bottom = 0.18 + title.h
        for section in subject.sections:
            name, content = next(iter(section.items()))
            content = [next(iter(item.items())) for item in content]
            content = {k: v for k, v in content}

            if "text" in content:

                text = Text(title=name, content=content["text"])
                text.add_to_shapes(
                    shapes=shapes,
                    left=0.25, top=_bottom+0.18
                )
                _bottom += 0.18 + text.h

                if "picture" in content:
                    fig = Figure(
                        title=content["picture"]["name"],
                        description=content["picture"]["description"],
                        pic_path=content["picture"]["path"],
                        size="small"
                    )

                    fig.add_to_shapes(
                        shapes=shapes,
                        left=0.25 + 0.62, top=_bottom+0.1
                    )

                    _bottom += 0.1 + fig.h

            elif "picture" in content:

                fig = Figure(
                        title=content["picture"]["name"],
                        description=content["picture"]["description"],
                        pic_path=content["picture"]["path"],
                        size="medium"
                    )

                fig.add_to_shapes(
                    shapes=shapes,
                    left=0.25, top=_bottom+0.18
                )

                _bottom += 0.18 + fig.h

        # add banners
        shapes.add_picture(
            str(PROJECT_DIR.joinpath('./data/banner_utechzone_blue.png')),
            left=0, top=self.prs.slide_height - Inch(1.07),
            width=self.prs.slide_width
            )


if __name__ == "__main__":

    yml = sys.argv[1]
    out_f = sys.argv[2]

    if not os.path.isabs(yml):
        settings = ReportSettings.from_yaml(
            yml,
            os.getcwd()
        )
    else:
        settings = ReportSettings.from_yaml(yml)

    if not os.path.isabs(out_f):
        out_f = pathlib.Path(os.getcwd()).joinpath(out_f)
        out_f = str(out_f)

    presentation = UTSimple(settings=settings)
    presentation.to_pptx(out_f)
