import sys
import os
import os.path as path
import string
import pptx
from pptx.util import Inches as Inch
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR as Line

_path = os.path.dirname(__file__)
if _path not in sys.path:
    sys.path.append(_path)

from _utils import line_counter  # noqa: E402


SUBJECT_TITLE_LIMITS = 35
DESCRIPTION_LINE_CAPACITY = 83
TEXT_TITLE_LIMITS = 57
TEXT_LINE_CAPACITY = 83


class SubjectTitle:

    def __init__(self, title, description=None):
        self._title = ""
        self.title = title
        self._title_font = "Arial"
        self._des = None
        self._des_lines = 0
        self._des_font = "Times New Roman"

        self._w = 10.25
        self._h = 0.6
        if description is not None:
            self.description = description

    @property
    def title(self):
        return self._title

    @title.setter
    def title(self, value):
        value = str(value)
        if len(value) > SUBJECT_TITLE_LIMITS:
            msg = "Title can has most {} characters"
            raise ValueError(msg.format(SUBJECT_TITLE_LIMITS))

        # make first word of title upper case
        value = string.capwords(value)
        self._title = value

    @property
    def description(self):
        return self._des

    @description.setter
    def description(self, value):
        self._des = str(value)
        self._des_lines = line_counter(
            self._des,
            DESCRIPTION_LINE_CAPACITY
        )
        self._h = 0.6 + 0.2 + self._des_lines*0.27

    @property
    def w(self):
        return self._w

    @property
    def h(self):
        return self._h

    def add_to_shapes(
            self,
            shapes,
            left=0.0, top=0.0
            ):
        """Add SubjectTitle into given shapes

        Args:
            shapes: the slide refernce to add
            left, top: specify the top-left corner of the object in Inch
        """
        left = Inch(left)
        top = Inch(top)

        # add subject title
        title = shapes.add_textbox(
            left=left, top=top,
            height=Inch(0.6), width=Inch(self.w)
            )
        param = title.text_frame.paragraphs[0]
        param.alignment = PP_ALIGN.LEFT

        run = param.add_run()
        run.text = self._title
        font = run.font
        font.name = self._title_font
        font.size = Pt(28)
        font.color.rgb = RGBColor(38, 38, 38)

        # draw a line under title
        shapes.add_connector(
            Line.STRAIGHT,
            # x, y, x_end, y_end
            left + Inch(0.08),
            top + Inch(0.57),
            left + Inch(0.08) + Inch(0.2*len(self._title)),
            top + Inch(0.57)
            )

        # add description
        textbox = shapes.add_textbox(
            left=left + Inch(0.91),
            top=top + Inch(0.6) + Inch(0.2),
            height=Inch(self.h) - Inch(0.8),
            width=Inch(self.w) - Inch(0.91)
            )
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        param = textbox.text_frame.paragraphs[0]
        param.alignment = PP_ALIGN.LEFT
        run = param.add_run()
        run.text = self._des
        font = run.font
        font.name = self._des_font
        font.size = Pt(16)
        font.color.rgb = RGBColor(64, 64, 64)


class Text:

    def __init__(self, title, content=None):
        self._title = ""
        self._title_font = "Times New Roman"
        self._content = None
        self._content_lines = 0
        self._content_font = "Times New Roman"

        self.title = title
        self._w = 10.17
        self._h = 0.27
        if content is not None:
            self.content = content

    @property
    def title(self):
        return self._title

    @title.setter
    def title(self, value):
        value = str(value)
        if len(value) > TEXT_TITLE_LIMITS:
            msg = "Title can has most {} characters"
            raise ValueError(msg.format(TEXT_TITLE_LIMITS))

        # make first word of title upper case
        value = string.capwords(value)
        self._title = value

    @property
    def content(self):
        return self._content

    @content.setter
    def content(self, value):
        self._content = str(value)
        self._content_lines = line_counter(
            self._content,
            TEXT_LINE_CAPACITY
        )
        self._h = 0.27 + self._content_lines*0.27

    @property
    def w(self):
        return self._w

    @property
    def h(self):
        return self._h

    def add_to_shapes(
            self,
            shapes,
            left=0.0, top=0.0
            ):
        """Add SubjectTitle into given shapes

        Args:
            shapes: the shape refernce to add
            left, top: specify the top-left corner of the object in Inch
        """
        left = Inch(left)
        top = Inch(top)

        # add text title
        textbox = shapes.add_textbox(
            left=left, top=top,
            height=Inch(self.h), width=Inch(self.w)
            )
        param = textbox.text_frame.paragraphs[0]
        param.alignment = PP_ALIGN.LEFT

        run = param.add_run()
        run.text = "#" + self.title
        font = run.font
        font.name = self._title_font
        font.size = Pt(16)
        font.bold = True
        font.color.rgb = RGBColor(64, 64, 64)

        # add text content
        param = textbox.text_frame.add_paragraph()
        param.alignment = PP_ALIGN.LEFT
        param.level = 1
        run = param.add_run()
        run.text = self.content
        font = run.font
        font.name = self._content_font
        font.size = Pt(16)
        font.color.rgb = RGBColor(64, 64, 64)
