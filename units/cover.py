import string
import datetime
import sys
import os.path as path

import pptx
from pptx.util import Inches as Inch
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR as Line


REPORT_TITLE_LIMITS = 20
SUBJECT_NUM_LIMITS = 3
SUBJECT_TITLE_LIMITS = 35


def week_regime(date):
    last_monday = date - datetime.timedelta(days=date.weekday())
    last_friday = last_monday + datetime.timedelta(days=4)
    return last_monday, last_friday


class ReportCover:

    def __init__(self, header, titles, author, date):
        self._w = 6.75
        self._h = 3.23

        self._titles = []
        self._titles_font = "Times New Roman"
        self.titles = titles

        self._header = ""
        self._header_font = "Arial"
        self.header = header

        self.author = author
        start, end = week_regime(date)
        self.date = (start, end)

    @property
    def titles(self):
        return self._titles

    @titles.setter
    def titles(self, titles):
        if isinstance(titles, str):
            titles = [string.capwords(titles)]
        else:
            titles = [string.capwords(title) for title in titles]

        if len(self.titles) > SUBJECT_NUM_LIMITS:
            msg = "Cover: Can have at most {} subjects"
            raise ValueError(msg.format(SUBJECT_NUM_LIMITS))
        for text in self.titles:
            if len(text) > SUBJECT_TITLE_LIMITS:
                msg = "Subject can has at most {} characters."
                raise ValueError(msg.format(SUBJECT_TITLE_LIMITS))

        self._titles.extend(titles)

    @property
    def header(self):
        return self._header

    @header.setter
    def header(self, value):
        value = str(value).upper()
        if len(value) > REPORT_TITLE_LIMITS:
            msg = "Report title can has at most {} characters."
            raise ValueError(msg.format(REPORT_TITLE_LIMITS))
        self._header = value

    @property
    def w(self):
        return self._w

    @property
    def h(self):
        return self._h

    def add_to_shapes(
            self,
            shapes,
            left=0.0, top=0.0,
            add_author_and_date=True
            ):
        """Add ReportCover into given shapes

        Args:
            shapes: the shapes refernce to add
            left, top: specify the top-left corner of the object in Inch
        """
        left = Inch(left)
        top = Inch(top)

        textbox = shapes.add_textbox(
            left=left, top=top,
            height=Inch(self.h), width=Inch(self.w)
            )
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        param = textbox.text_frame.paragraphs[0]
        param.alignment = PP_ALIGN.CENTER

        # Cover Text
        run = param.add_run()
        run.text = self._header
        font = run.font
        font.name = self._header_font
        font.size = Pt(54)
        font.color.rgb = RGBColor(127, 127, 127)

        # Draw a line under title
        shapes.add_connector(
            Line.STRAIGHT,
            # x, y, x_end, y_end7
            left+Inch(0.32), top+Inch(0.93), left+Inch(6.57), top+Inch(0.93)
            )

        for text in self.titles:
            param = textbox.text_frame.add_paragraph()
            param.space_before = Pt(6)
            param.alignment = PP_ALIGN.CENTER
            run = param.add_run()
            run.text = text
            font = run.font
            font.name = self._titles_font
            font.size = Pt(44)
            font.color.rgb = RGBColor(127, 127, 127)

        # Add author information
        if add_author_and_date:
            info = "Author - {}; Date - {} ~ {}"
            textbox = shapes.add_textbox(
                left=left - Inch(0.24),
                top=top + Inch(3.91),
                height=Inch(0.4), width=Inch(7.65)
                )
            param = textbox.text_frame.paragraphs[0]
            param.alignment = PP_ALIGN.CENTER

            run = param.add_run()
            run.text = info.format(
                self.author,
                self.date[0],
                self.date[1]
                )
            font = run.font
            font.name = self._header_font
            font.size = Pt(18)
            font.color.rgb = RGBColor(84, 142, 213)
