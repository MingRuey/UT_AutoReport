import os.path as path
import cv2
import pptx
from pptx.util import Inches as Inch
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

DATADIR = path.join(path.dirname(path.abspath(__file__)), "data")
PIXEL_TO_INCH = 1/96


def _get_image_shape(img_path):
    img = cv2.imread(str(img_path))
    if img is None:
        msg = "Failed to load image: {}"
        raise ValueError(msg.format(img_path))

    height, width, _channel = img.shape
    height = height * PIXEL_TO_INCH
    width = width * PIXEL_TO_INCH
    return height, width


def _get_resize_shape(raw_shape, target_shape):
    raw_h, raw_w = raw_shape
    target_h, target_w = target_shape

    if raw_h/target_h >= raw_w/target_w:
        new_h = target_h
        new_w = raw_w * target_h / raw_h
    else:
        new_w = target_w
        new_h = raw_h * target_w / raw_w
    return int(new_h), int(new_w)


class Figure:

    SMALL_SHAPE = (2.81, 3.54)
    MEDIUM_SHAPE = (3.96, 5)
    BIG_SHAPE = (6.83, 4.96)

    def __init__(self, title, description, pic_path, size=None):

        size = str(size).lower()
        if size == "small":
            self._size = self.SMALL_SHAPE
        elif size == "medium":
            self._size = self.MEDIUM_SHAPE
        elif size == "big":
            self._size = self.BIG_SHAPE
        else:
            msg = "Not supported size option: {}"
            raise NotImplementedError(msg.format(size))

        self._title = str(title)
        self._des = str(description)
        self._font = "Times New Roman"
        self._font_size = Pt(14) if size == "small" else Pt(16)
        self._pic = pic_path

        self._pic_h, self._pic_w = _get_resize_shape(
            raw_shape=_get_image_shape(pic_path),
            target_shape=self._size
        )
        self._h = self.pic_h
        self._w = 10.58

    @property
    def pic_path(self):
        return self._pic

    @property
    def pic_h(self):
        return self._pic_h

    @property
    def pic_w(self):
        return self._pic_w

    @property
    def h(self):
        return self._h

    @property
    def w(self):
        return self._w

    @property
    def title(self):
        return self._title

    @property
    def description(self):
        return self._des

    def add_to_shapes(self, shapes, left=0.0, top=0.0):
        """Add Figure into given shapes

        Args:
            shapes: the shape refernce to add
            left, top: specify the top-left corner of the object in Inch
        """
        left = Inch(left)
        top = Inch(top)

        shapes.add_picture(
            image_file=self.pic_path,
            left=left, top=top,
            width=Inch(self.pic_w), height=Inch(self.pic_h)
            )

        from_pic_to_text = Inch(0.17)
        textbox = shapes.add_textbox(
            left=left + Inch(self.pic_w) + from_pic_to_text,
            top=top,
            height=Inch(self.h),
            width=Inch(self.w) - Inch(self.pic_w) - from_pic_to_text - left
            )
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        param = textbox.text_frame.paragraphs[0]
        param.alignment = PP_ALIGN.LEFT

        # Cover Text
        run = param.add_run()
        run.text = self._title + "\n\n" + self._des
        font = run.font
        font.name = self._font
        font.size = self._font_size
        font.color.rgb = RGBColor(38, 38, 38)


if __name__ == "__main__":
    pass
